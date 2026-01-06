#!/usr/bin/env python3
"""
Generate youlidao.ai Introduction PowerPoint
Based on youlidao-ai-ppt-outline.md
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
PSYCHIC_DARK = RGBColor(0x4A, 0x14, 0x8C)  # Deep purple
NEON_CYAN = RGBColor(0x00, 0xE5, 0xFF)      # Highlight
INDUSTRIAL_GRAY = RGBColor(0x37, 0x47, 0x4F) # Background
NEON_GREEN = RGBColor(0x00, 0xE6, 0x76)     # Success
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)

def set_slide_background(slide, color):
    """Set solid background color for a slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, INDUSTRIAL_GRAY)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = NEON_CYAN
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets, highlight_box=None):
    """Add a content slide with bullets"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, INDUSTRIAL_GRAY)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Accent line under title
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(2), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = NEON_CYAN
    line.line.fill.background()

    # Bullets
    bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4.5))
    tf = bullet_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(22)
        p.font.color.rgb = LIGHT_GRAY
        p.space_after = Pt(12)
        p.level = 0

    # Highlight box if provided
    if highlight_box:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.5), Inches(5.5), Inches(9), Inches(0.8))
        box.fill.solid()
        box.fill.fore_color.rgb = PSYCHIC_DARK
        box.line.color.rgb = NEON_CYAN
        box.line.width = Pt(2)

        tf = box.text_frame
        tf.paragraphs[0].text = highlight_box
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = NEON_CYAN
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.italic = True

    return slide

def add_comparison_slide(prs, title, headers, rows):
    """Add a comparison table slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, INDUSTRIAL_GRAY)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Table
    cols = len(headers)
    table_rows = len(rows) + 1  # +1 for header

    table = slide.shapes.add_table(table_rows, cols, Inches(0.5), Inches(1.3),
                                    Inches(9), Inches(4)).table

    # Set column widths
    col_width = Inches(9 / cols)
    for i in range(cols):
        table.columns[i].width = col_width

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = PSYCHIC_DARK
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = INDUSTRIAL_GRAY
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = LIGHT_GRAY
            p.alignment = PP_ALIGN.CENTER

    return slide

def add_team_slide(prs, title, team_members):
    """Add a team introduction slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, INDUSTRIAL_GRAY)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Team cards
    card_width = Inches(1.6)
    card_height = Inches(1.8)
    start_x = Inches(0.5)
    start_y = Inches(1.5)
    gap = Inches(0.2)

    for i, (icon, role, desc) in enumerate(team_members):
        col = i % 5
        row = i // 5
        x = start_x + col * (card_width + gap)
        y = start_y + row * (card_height + gap)

        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = PSYCHIC_DARK
        card.line.color.rgb = NEON_CYAN
        card.line.width = Pt(1)

        # Icon
        icon_box = slide.shapes.add_textbox(x, y + Inches(0.1), card_width, Inches(0.5))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(28)
        p.alignment = PP_ALIGN.CENTER

        # Role name
        role_box = slide.shapes.add_textbox(x, y + Inches(0.6), card_width, Inches(0.4))
        tf = role_box.text_frame
        p = tf.paragraphs[0]
        p.text = role
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = NEON_CYAN
        p.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide.shapes.add_textbox(x, y + Inches(1.0), card_width, Inches(0.7))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_section_title(prs, title, subtitle=""):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, PSYCHIC_DARK)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = NEON_CYAN
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_pricing_slide(prs):
    """Add pricing slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, INDUSTRIAL_GRAY)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "简单透明的定价"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Pricing cards
    plans = [
        ("访客计划", "$399/月", "初次体验"),
        ("常住计划", "$369/月", "持续开发"),
        ("创始人计划", "$299/月", "早期支持者特惠")
    ]

    card_width = Inches(2.8)
    card_height = Inches(2.5)
    start_x = Inches(0.5)
    gap = Inches(0.3)

    for i, (name, price, desc) in enumerate(plans):
        x = start_x + i * (card_width + gap)
        y = Inches(1.5)

        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = PSYCHIC_DARK

        # Highlight the founder plan
        if i == 2:
            card.line.color.rgb = NEON_GREEN
            card.line.width = Pt(3)
        else:
            card.line.color.rgb = NEON_CYAN
            card.line.width = Pt(1)

        # Plan name
        name_box = slide.shapes.add_textbox(x, y + Inches(0.2), card_width, Inches(0.4))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Price
        price_box = slide.shapes.add_textbox(x, y + Inches(0.7), card_width, Inches(0.6))
        tf = price_box.text_frame
        p = tf.paragraphs[0]
        p.text = price
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = NEON_CYAN if i != 2 else NEON_GREEN
        p.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide.shapes.add_textbox(x, y + Inches(1.4), card_width, Inches(0.5))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER

    # Benefits list
    benefits = [
        "✅ 完整 9 人 AI 团队",
        "✅ 无限项目数量",
        "✅ 无限命令执行",
        "✅ Git 深度集成",
        "✅ 多设备同步"
    ]

    benefit_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(9), Inches(2))
    tf = benefit_box.text_frame

    for i, benefit in enumerate(benefits):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = benefit
        p.font.size = Pt(16)
        p.font.color.rgb = LIGHT_GRAY

    return slide

def add_cta_slide(prs):
    """Add call to action slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, INDUSTRIAL_GRAY)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "三步开启 AI 软件公司"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Steps
    steps = [
        ("1", "注册", "访问 youlidao.ai"),
        ("2", "订阅", "选择适合的方案"),
        ("3", "创建", "开始与 AI 团队协作")
    ]

    step_width = Inches(2.5)
    start_x = Inches(1)
    gap = Inches(0.5)
    y = Inches(2)

    for i, (num, action, desc) in enumerate(steps):
        x = start_x + i * (step_width + gap)

        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.75), y, Inches(1), Inches(1))
        circle.fill.solid()
        circle.fill.fore_color.rgb = NEON_CYAN
        circle.line.fill.background()

        # Number
        num_box = slide.shapes.add_textbox(x + Inches(0.75), y + Inches(0.15), Inches(1), Inches(0.7))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = INDUSTRIAL_GRAY
        p.alignment = PP_ALIGN.CENTER

        # Action
        action_box = slide.shapes.add_textbox(x, y + Inches(1.2), step_width, Inches(0.5))
        tf = action_box.text_frame
        p = tf.paragraphs[0]
        p.text = action
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide.shapes.add_textbox(x, y + Inches(1.7), step_width, Inches(0.5))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER

        # Arrow between steps
        if i < 2:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                           x + step_width + Inches(0.1), y + Inches(0.4),
                                           Inches(0.3), Inches(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = NEON_CYAN
            arrow.line.fill.background()

    # CTA Button
    btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(3), Inches(4.5), Inches(4), Inches(0.8))
    btn.fill.solid()
    btn.fill.fore_color.rgb = NEON_CYAN
    btn.line.fill.background()

    btn_text = slide.shapes.add_textbox(Inches(3), Inches(4.6), Inches(4), Inches(0.6))
    tf = btn_text.text_frame
    p = tf.paragraphs[0]
    p.text = "🚀 立即体验"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = INDUSTRIAL_GRAY
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_end_slide(prs):
    """Add ending slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, PSYCHIC_DARK)

    # Main slogan
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "你的 AI 软件公司"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "随时待命"
    p.font.size = Pt(36)
    p.font.color.rgb = NEON_CYAN
    p.alignment = PP_ALIGN.CENTER

    # Vision
    vision_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
    tf = vision_box.text_frame
    p = tf.paragraphs[0]
    p.text = "让每一个有想法的人，都能把创意变成现实"
    p.font.size = Pt(20)
    p.font.italic = True
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(1.2))
    tf = contact_box.text_frame

    p = tf.paragraphs[0]
    p.text = "🌐 youlidao.ai"
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "📧 contact@youlidao.ai"
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    return slide

def main():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # === Slide 1: Cover ===
    add_title_slide(prs, "AI 驱动的软件公司", "9 位 AI 工程师，实时为你开发软件")

    # === Slide 2: Pain Points ===
    add_content_slide(prs, "开发软件，为什么这么难？", [
        "👥 人才贵 — 一个 5 人团队年薪成本 > ¥200 万",
        "⏰ 周期长 — 一个 MVP 平均需要 3-6 个月",
        "🔒 黑盒化 — 外包开发进度不透明，沟通成本高"
    ], "有没有一种方式，既快又透明？")

    # === Slide 3: Market Opportunity ===
    add_content_slide(prs, "AI 正在重塑软件行业", [
        "📈 2024 年 AI 编程工具市场规模持续高速增长",
        "👨‍💻 GitHub Copilot 用户突破 100 万+",
        "⚠️ 但现有工具只解决了「编码」这一个环节",
        "",
        "🎯 市场上缺少覆盖完整开发流程的 AI 解决方案"
    ])

    # === Slide 4: Core Concept ===
    add_comparison_slide(prs, "透明的 AI 软件公司",
        ["对比维度", "传统 AI 工具", "youlidao.ai"],
        [
            ["团队规模", "单一助手", "完整 9 人团队"],
            ["覆盖范围", "只管编码", "需求→上线全流程"],
            ["透明度", "黑盒输出", "透明厨房，全程可见"],
            ["质量保障", "无", "内置质量门禁"]
        ])

    # === Slide 5: Strategic Team ===
    add_team_slide(prs, "战略规划部 — 5 位专家", [
        ("🔍", "Analyst", "市场调研\n项目初探"),
        ("📋", "PM", "需求梳理\n产品规划"),
        ("🎨", "UX-Expert", "界面设计\n交互规范"),
        ("🏗️", "Architect", "技术选型\n系统架构"),
        ("✅", "PO", "质量把关\n优先级管理")
    ])

    # === Slide 6: Engineering Team ===
    add_team_slide(prs, "工程实施部 — 4 位专家", [
        ("📊", "SM", "故事拆分\n迭代管理"),
        ("🔧", "Architect-Eng", "技术评审\n代码规范"),
        ("💻", "Dev", "代码实现\n功能开发"),
        ("🧪", "QA", "测试验证\n质量门禁")
    ])

    # === Slide 7: Workflow ===
    add_content_slide(prs, "从想法到上线，一条龙服务", [
        "📥 需求输入 → Analyst 分析 → PM 规划 → UX 设计",
        "",
        "⬇️ 技术方案 → Architect 架构设计",
        "",
        "⬇️ 开发实施 → SM 拆分 → Dev 实现 → QA 验证",
        "",
        "📤 上线交付"
    ], "自动流转 | 质量门禁 | 全程可追溯")

    # === Slide 8: Product Interface ===
    add_content_slide(prs, "三区工作台 — 全景掌控", [
        "📁 左侧：档案室 — 实时查看每一个文件变化",
        "",
        "👥 中上：工程师面板 — 9 位 AI 实时状态指示灯",
        "",
        "🖥️ 中下：工程监视器 — 实时观看 AI 工作过程",
        "",
        "💬 右侧：指令通道 — 自然语言下达命令"
    ])

    # === Slide 9: Real-time Status ===
    add_content_slide(prs, "透明厨房 — 每一步都看得见", [
        "🔘 空闲 (Idle) — 等待任务分配",
        "",
        "🟡 激活中 (Activating) — 正在启动工作环境",
        "",
        "🟢 就绪 (Ready) — 可以接受命令",
        "",
        "🔵 执行中 (Executing) — 正在努力工作"
    ], "告别黑盒，每个 AI 的工作状态一目了然")

    # === Slide 10: Git Integration ===
    add_content_slide(prs, "与你的代码库无缝对接", [
        "1️⃣ 粘贴 GitHub URL — 支持任何 Git 仓库",
        "",
        "2️⃣ 自动配置 SSH 密钥 — 一键授权，安全可靠",
        "",
        "3️⃣ AI 团队开始工作 — 理解现有架构后再迭代"
    ], "支持导入现有项目，在你的代码基础上继续开发")

    # === Section: Use Cases ===
    add_section_title(prs, "应用场景", "谁在使用 AI 软件公司？")

    # === Slide 11: Startup ===
    add_comparison_slide(prs, "场景一：创业者的 MVP 加速器",
        ["对比项", "传统方式", "AI 软件公司"],
        [
            ["启动时间", "招聘团队 2-3 月", "即刻开始"],
            ["沟通成本", "外包沟通效率低", "直接对话 AI"],
            ["预算需求", "¥50 万+", "¥3000/月起"],
            ["交付周期", "数月", "数天到原型"]
        ])

    # === Slide 12: Tech Team ===
    add_content_slide(prs, "场景二：技术团队的效能倍增器", [
        "😫 痛点：需求积压、人手不足、重复劳动多",
        "",
        "🔄 接手重复性开发任务",
        "📝 自动生成样板代码",
        "🧪 自动化测试编写",
        "📚 文档自动生成"
    ], "团队产能提升 3-5 倍")

    # === Slide 13: PM ===
    add_content_slide(prs, "场景三：产品经理的快速验证工具", [
        "💡 几小时内产出可交互原型",
        "",
        "📊 技术可行性即时评估",
        "",
        "🔀 多方案快速对比",
        "",
        "💬 \"以前需要等开发排期，现在下午就能看到效果\""
    ])

    # === Section: Why Us ===
    add_section_title(prs, "为什么选择我们", "不止是 AI 编程助手")

    # === Slide 14: Comparison ===
    add_comparison_slide(prs, "竞品对比",
        ["特性", "youlidao", "Cursor", "Copilot", "Replit"],
        [
            ["完整 AI 团队", "✅", "❌", "❌", "❌"],
            ["全流程覆盖", "✅", "❌", "❌", "⚠️"],
            ["透明可见", "✅", "⚠️", "❌", "⚠️"],
            ["架构设计", "✅", "❌", "❌", "❌"],
            ["质量门禁", "✅", "❌", "❌", "⚠️"]
        ])

    # === Slide 15: Technology ===
    add_content_slide(prs, "企业级技术架构", [
        "🧠 Claude AI — 业界最强推理能力的 AI 引擎",
        "",
        "⚡ 实时同步 — WebSocket 支持多设备协作",
        "",
        "🔒 银行级安全 — 行级数据隔离，代码只属于你",
        "",
        "🌍 全球 CDN — Vercel Edge 极速响应"
    ])

    # === Slide 16: Pricing ===
    add_pricing_slide(prs)

    # === Slide 17: CTA ===
    add_cta_slide(prs)

    # === Slide 18: End ===
    add_end_slide(prs)

    # Save
    output_path = "youlidao-ai-intro.pptx"
    prs.save(output_path)
    print(f"✅ PPT 已生成: {output_path}")
    print(f"📊 共 {len(prs.slides)} 页")

if __name__ == "__main__":
    main()
